from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from app.services.parsers.base import ParsedSegmentPayload
from app.utils.text import clean_text, split_text_blocks


def parse_pptx_file(path: Path) -> list[ParsedSegmentPayload]:
    presentation = Presentation(path)
    segments: list[ParsedSegmentPayload] = []
    order = 1
    for slide_index, slide in enumerate(presentation.slides, start=1):
        text_parts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text_parts.append(shape.text)
        slide_text = clean_text("\n".join(text_parts))
        if not slide_text:
            continue
        blocks = split_text_blocks(slide_text, min_length=18)
        for block in blocks:
            text = clean_text(block)
            if not text:
                continue
            segments.append(
                ParsedSegmentPayload(
                    source_type="pptx",
                    source_label=f"PPT 第 {slide_index} 页",
                    segment_order=order,
                    content=text,
                    metadata_json={"slide": slide_index},
                )
            )
            order += 1
    return segments
